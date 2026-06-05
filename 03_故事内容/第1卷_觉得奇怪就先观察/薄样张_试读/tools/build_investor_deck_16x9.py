#!/usr/bin/env python3
"""Thin investor deck · 16:9 · 5–7 slides · trial-read focused (not full expert pack).

Output:
  薄样张_试读/PDF/学堂奇事録_Vol1_薄样张_路演_{date}.pptx
"""

from __future__ import annotations

import subprocess
import sys
from datetime import date
from pathlib import Path

TRIAL = Path(__file__).resolve().parents[1]
PDF_DIR = TRIAL / "PDF"
OUT_DATE = date.today().strftime("%Y%m%d")
PPT_NAME = f"学堂奇事録_Vol1_薄样张_路演_{OUT_DATE}.pptx"


def build_ppt() -> Path:
    try:
        from pptx import Presentation
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt
    except ImportError:
        subprocess.run([sys.executable, "-m", "pip", "install", "python-pptx", "-q"], check=True)
        from pptx import Presentation
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt

    W16 = Inches(16)
    H9 = Inches(9)
    MX = Inches(0.65)
    MY = Inches(0.45)
    CW = W16 - 2 * MX

    prs = Presentation()
    prs.slide_width = W16
    prs.slide_height = H9
    blank = prs.slide_layouts[6]

    def title_slide(title: str, subtitle: str = "", tagline: str = "") -> None:
        s = prs.slides.add_slide(blank)
        box = s.shapes.add_textbox(MX, Inches(2.0), CW, Inches(4.5))
        tf = box.text_frame
        tf.text = title
        tf.paragraphs[0].font.size = Pt(40)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        if tagline:
            p = tf.add_paragraph()
            p.text = tagline
            p.font.size = Pt(22)
            p.font.italic = True
            p.alignment = PP_ALIGN.CENTER
        if subtitle:
            p = tf.add_paragraph()
            p.text = subtitle
            p.font.size = Pt(18)
            p.alignment = PP_ALIGN.CENTER
        foot = s.shapes.add_textbox(MX, Inches(8.35), CW, Inches(0.4))
        foot.text_frame.text = f"薄样张路演 · 16:9 · {OUT_DATE}"
        foot.text_frame.paragraphs[0].font.size = Pt(10)
        foot.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def bullets(title: str, lines: list[str]) -> None:
        s = prs.slides.add_slide(blank)
        tb = s.shapes.add_textbox(MX, MY, CW, Inches(0.7))
        tb.text_frame.text = title
        tb.text_frame.paragraphs[0].font.size = Pt(30)
        tb.text_frame.paragraphs[0].font.bold = True
        body = s.shapes.add_textbox(MX, Inches(1.15), CW, Inches(7.2))
        tf = body.text_frame
        tf.word_wrap = True
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.size = Pt(17)
            p.space_after = Pt(8)

    title_slide(
        "学堂奇事録",
        "第1巻 · おかしいと思ったら、まず見てみる",
        "謎が、解ける日。",
    )

    bullets(
        "首市场 · 日本",
        [
            "対象：7–11歳（主 8–10）· 名古屋 · 校园軽推理成長 IP",
            "境界：硬科普/参考書/恐怖/大人捜査 ✗",
            "設計：かわいさで入る · 謎で留まる · 観察で続く",
            "Vol1 = Plan B · 5年跨班 観察クラブ · 5小事件合订",
            "本轮：案①薄样张试读 → 5–8組 E20 后再扩写",
        ],
    )

    bullets(
        "差別化 · 観察クラブ",
        [
            "学校おもしろ観察クラブ — 怪しいより「おかしい」",
            "陸珣（しゅん）視点 · 転校生 · 微观观察 · 非少年侦探团",
            "陸瑆（ひかる · 4年2組）笔记层 · 温柔一问收束",
            "公平线索 · 温柔真相 · 零恐怖 · 名古屋上履き",
            "机制谜底不可被人际替换（正典门禁）",
        ],
    )

    bullets(
        "案① · めくれたポスター（试读核心）",
        [
            "奇怪：壁报海报的边，每天翘在不同侧",
            "线索：风侧 · 湿度 · 贴法方向 — 读者可猜",
            "功能：入社前「先观察」· 侧廊日常 · 插图可读",
            "尾钩：壁报留空一栏 → 串联 Vol1 五案",
            "试读 PDF：序+案①+瑆页+实验+问卷（~20–30p）",
        ],
    )

    bullets(
        "试读计划 · E20",
        [
            "对象：5–8 组 · 8–10 岁 + 家长",
            "交付：日本語薄样张 PDF + 主持人脚本 + 问卷",
            "阈值：「続きが読みたい」≥60% · scary 0% · 安全 100%",
            "未达标：改样章/插图/尾钩 — 不先扩写案②–⑤",
            "通过后：Nano Banana 8角色终稿 + 全卷排版",
        ],
    )

    bullets(
        "开放门禁 · 尽调一句话",
        [
            "✅ 正典：Vol1 Plan B · 首市场日本 · 名称 LOCK",
            "🟡 日文：A001 JP_VOICE_v1 · ②–⑤初译待 E04",
            "🟡 插图：V-S01 可用 · 角色 L0 待画师/Nano 终稿",
            "⬜ 市场：E20 试读数据 = 0 → 本轮首要风险",
            "内容脚手架 84 分 · 市场验证 0 分 → 先试读再出版包",
        ],
    )

    title_slide(
        "Next · 试读反馈后",
        "全卷5案 · 出版成果 PDF/PPT · 渠道洽谈",
        "薄样张 ≠ 正式版 05_出版成果",
    )

    PDF_DIR.mkdir(parents=True, exist_ok=True)
    out = PDF_DIR / PPT_NAME
    prs.save(str(out))
    return out


def main() -> int:
    path = build_ppt()
    print(f"PPTX: {path} ({path.stat().st_size // 1024} KB)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
