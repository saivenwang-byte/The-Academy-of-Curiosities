#!/usr/bin/env python3
"""MVP V2 Unit1 A001-A005 — deadline delivery assembler.

Outputs under 05_出版成果/MVP_V2_20260608/:
  00_MVP交付清单.md
  01_五案CN/  02_五案JP/  03_插图/
  reader.html
  学堂趣事录_第1单元_MVP试读_V2.0_20260608.pdf
  第1单元_MVP路演_V2.0.pptx
"""

from __future__ import annotations

import argparse
import base64
import re
import shutil
import subprocess
import sys
from datetime import date
from pathlib import Path

TOOLS = Path(__file__).resolve().parent
FORMAL = TOOLS.parents[1]
ROOT = FORMAL.parents[3]
BODY = FORMAL / "01_正文"
DEPTH = FORMAL.parent / "样章包" / "插图" / "depth_anchor"
DISPLAY_TOOLS = FORMAL / "04_展示版" / "tools"
OUT_ROOT = FORMAL / "05_出版成果" / "MVP_V2_20260608"
MIGRATION = FORMAL.parent / "V2迁移"

OUT_DATE = "20260608"
PDF_NAME = f"学堂趣事录_第1单元_MVP试读_V2.0_{OUT_DATE}.pdf"
PDF_FULL_CN_NAME = "学堂趣事录_第1单元_读者试读_FULL_V2.0.pdf"
PPT_NAME = f"第1单元_MVP路演_V2.0_{OUT_DATE}.pptx"

CASES = [
    {
        "id": "A001",
        "num": "一",
        "cn_title": "全班都听见了他的声音",
        "jp_title": "クラス全員が彼の声を聞いた",
        "cn_file": "案01_全班都听见了他的声音_HybridVoice_V2.0.txt",
        "jp_file": "案01_全班都听见了他的声音_HybridVoice_V2.0_日本語.txt",
        "illus": [
            ("V-S01-V2-A1_广播响起_G1draft_c02.png", "A001 · 放送 · 唇が合わない · G1 R2"),
            ("V-S01-V2-A3_文件时间_G1draft_c02.png", "A001 · 0328 · 旧録音 · G1 R2"),
        ],
        "illus_subdir": "案01",
        "hook": "放送の声は光のものではない——ファイルと波形が先に語る。",
    },
    {
        "id": "A002",
        "num": "二",
        "cn_title": "没有人写过的道歉",
        "jp_title": "誰も書いていない謝罪",
        "cn_file": "案02_没有人写过的道歉_HybridVoice_V2.0.txt",
        "jp_file": "案02_没有人写过的道歉_HybridVoice_V2.0_日本語.txt",
        "illus": [
            ("V-S02-V2-A1_黑板对不起_G1draft_c02.png", "A002 · 黒板 · SC-02 · G1 R2"),
            ("V-S02-V2-A3_膜边反光_G1draft_c02.png", "A002 · 膜の端 · SC-05 · G1 R2"),
        ],
        "illus_subdir": "案02",
        "hook": "黒板に謝罪の字——書いた人は教室に入っていない。",
    },
    {
        "id": "A003",
        "num": "三",
        "cn_title": "每个人都记得的海报",
        "jp_title": "みんな覚えているポスター",
        "cn_file": "案03_每个人都记得的海报_HybridVoice_V2.0.txt",
        "jp_file": "案03_每个人都记得的海报_HybridVoice_V2.0_日本語.txt",
        "illus": [
            ("V-S03-V2-DEMO_空海报位_G1draft_c02.png", "A003 · 空欄 · SC-01 · G1 R2"),
            ("V-S03-V2-A4_远标题连线_G1draft_c02.png", "A003 · 遠タイトル · SC-04 · G1 R2"),
        ],
        "illus_subdir": "案03",
        "hook": "三人が三種類のポスターを覚えている——実物はない。",
    },
    {
        "id": "A004",
        "num": "四",
        "cn_title": "只出现在她抽屉里的失物",
        "jp_title": "彼女の引き出しにだけ現れた落とし物",
        "cn_file": "案04_只出现在她抽屉里的失物_HybridVoice_V2.0.txt",
        "jp_file": "案04_只出现在她抽屉里的失物_HybridVoice_V2.0_日本語.txt",
        "illus": [
            ("V-S04-V2-A1_抽屉失物_G1draft_c02.png", "A004 · 引き出し · SC-02 · G1 R2"),
            ("V-S04-V2-A3_倾斜水泡_G1draft_c02.png", "A004 · 傾き水準器 · SC-05 · G1 R2"),
            ("V-S04-V2-A4_振动复现_G1draft_c02.png", "A004 · 振動再現 · SC-07 · G1 R2"),
        ],
        "illus_subdir": "案04",
        "hook": "鍵のない引き出しに物が現れる——振動と傾きが鍵。",
    },
    {
        "id": "A005",
        "num": "五",
        "cn_title": "午休后消失的影子",
        "jp_title": "昼休みのあと消えた影",
        "cn_file": "案05_午休后消失的影子_HybridVoice_V2.0.txt",
        "jp_file": "案05_午休后消失的影子_HybridVoice_V2.0_日本語.txt",
        "illus": [
            ("V-S05-V2-A1_仅水野无影_G1draft_c02.png", "A005 · 影なし · SC-01 · G1 R2"),
            ("V-S05-V2-A6_重拍有影_G1draft_c02.png", "A005 · 再撮影 · SC-08 · G1 R2"),
        ],
        "illus_subdir": "案05",
        "hook": "集合写真で一人だけ影がない——タイムラインのずれ。",
    },
]

CHROME_CANDIDATES = [
    Path(r"C:\Program Files\Microsoft\Edge\Application\msedge.exe"),
    Path(r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe"),
    Path(r"C:\Program Files\Google\Chrome\Application\chrome.exe"),
    Path(r"C:\Program Files (x86)\Google\Chrome\Application\chrome.exe"),
]

sys.path.insert(0, str(DISPLAY_TOOLS))
import build_display_edition as bde  # noqa: E402

META_CUT = re.compile(
    r"\n---\n\n(?:【读者可以试|【EDITOR|【V2\.0|【G-JP · EDITOR|【G-JP PILOT|【Hybrid Voice)"
)


def strip_v2_reader(raw: str) -> str:
    m = META_CUT.search(raw)
    if m:
        raw = raw[: m.start()]
    lines: list[str] = []
    skip_block = False
    for line in raw.splitlines():
        s = line.strip()
        if s.startswith("【") and ("EDITOR" in s or "PRODUCTION" in s or "G-JP" in s or "正式版" in s):
            skip_block = True
            continue
        if skip_block and (s.startswith("·") or s.startswith("  ·") or not s):
            continue
        if skip_block and s and not s.startswith("·"):
            skip_block = False
        if s.startswith("《学堂") or s.startswith("『学堂"):
            continue
        if re.match(r"^第1[卷巻]", s):
            continue
        if s.startswith("单元 ·"):
            continue
        if re.match(r"^=+$", s) or re.match(r"^-+$", s):
            continue
        if s.startswith("【V2.0 场级"):
            break
        lines.append(line)
    return "\n".join(lines).strip()


def img_uri(path: Path | None) -> str:
    if not path or not path.exists():
        return ""
    data = base64.b64encode(path.read_bytes()).decode("ascii")
    ext = path.suffix.lower().lstrip(".")
    mime = "png" if ext == "png" else ext
    return f"data:image/{mime};base64,{data}"


def figure_block(path: Path | None, caption: str) -> str:
    uri = img_uri(path)
    if uri:
        return (
            f'<figure class="fig"><img src="{uri}" alt="{caption}"/>'
            f'<figcaption>{caption}</figcaption></figure>'
        )
    return (
        f'<div class="placeholder">'
        f'<p>【插图占位】</p><p>{caption}</p></div>'
    )


def prose_to_html(text: str, lang: str) -> str:
    blocks: list[str] = []
    pending_h3: str | None = None
    paras: list[str] = []

    def flush() -> None:
        nonlocal pending_h3, paras
        if pending_h3 is None and not paras:
            return
        inner = []
        if pending_h3:
            inner.append(f'<h3 class="sec">{pending_h3}</h3>')
        for p in paras:
            for ln in p.split("\n"):
                ln = ln.strip()
                if ln:
                    inner.append(f"<p>{ln}</p>")
        blocks.append('<div class="scene">' + "\n".join(inner) + "</div>")
        pending_h3 = None
        paras = []

    for para in re.split(r"\n\s*\n", text.strip()):
        para = para.strip()
        if not para:
            continue
        lines = para.splitlines()
        head = lines[0].strip()
        if head.startswith("序"):
            flush()
            body = "\n".join(lines[1:]).strip()
            blocks.append(f'<h2 class="chapter">{head}</h2>')
            if body:
                blocks.append(f"<p>{body.replace(chr(10), '</p><p>')}</p>")
            continue
        if re.match(r"^[一二三四五]、", head):
            flush()
            blocks.append(f'<h2 class="case-title">{head}</h2>')
            rest = "\n".join(lines[1:]).strip()
            if rest:
                paras.append(rest)
            continue
        if re.match(r"^###\s+", head):
            flush()
            pending_h3 = re.sub(r"^###\s+", "", head)
            rest = "\n".join(lines[1:]).strip()
            if rest:
                paras.append(rest)
            continue
        if head.startswith("瑆笔记") or head.startswith("瑆ノート"):
            flush()
            blocks.append(f'<aside class="hikaru"><h3>{head}</h3>')
            for ln in lines[1:]:
                ln = ln.strip().lstrip("> ").strip()
                if ln:
                    blocks.append(f"<p>{ln}</p>")
            blocks.append("</aside>")
            continue
        if pending_h3 is not None:
            paras.append(para)
        else:
            blocks.append(f"<p>{para.replace(chr(10), '</p><p>')}</p>")
    flush()
    return "\n".join(blocks)


def build_reader_html(lang: str) -> str:
    is_ja = lang == "ja"
    title = "学堂奇事録" if is_ja else "学堂趣事录"
    vol = (
        "第1巻 · へんなところ、先に見てみる"
        if is_ja
        else "第1卷 · 觉得奇怪，就先观察"
    )
    badge = "MVP V2 · 2026-06-08 · 非出版清样"
    parts = [
        f"""<!DOCTYPE html>
<html lang="{'ja' if is_ja else 'zh-CN'}">
<head>
<meta charset="utf-8"/>
<title>{title} · 第1单元 MVP V2.0</title>
<style>
@page {{ size: A5; margin: 14mm; }}
body {{
  font-family: {"'Yu Gothic', 'Meiryo', sans-serif" if is_ja else "'Microsoft YaHei', 'PingFang SC', sans-serif"};
  font-size: 10.5pt; line-height: 1.55; color: #222; max-width: 128mm; margin: 0 auto;
}}
.cover {{ text-align: center; padding: 2em 0 3em; page-break-after: always; }}
.badge {{ background: #8b3a2a; color: #fff; font-size: 8pt; padding: 0.2em 0.5em; display: inline-block; }}
h1 {{ font-size: 22pt; margin: 0.4em 0; }}
h2.chapter, h2.case-title {{ font-size: 14pt; border-bottom: 1px solid #ccc; padding-bottom: 0.2em; margin-top: 1.2em; }}
h3.sec {{ font-size: 11pt; color: #444; }}
.case-block {{ page-break-before: always; }}
.fig img {{ max-width: 100%; height: auto; display: block; margin: 0.6em auto; }}
.fig figcaption, .placeholder {{ font-size: 8.5pt; color: #666; text-align: center; }}
.placeholder {{ border: 2px dashed #bbb; padding: 1.5em; margin: 0.8em 0; }}
.hikaru {{ background: #fffef5; border-left: 3px solid #d4c4a8; padding: 0.8em 1em; margin: 1em 0; font-size: 10pt; }}
p {{ margin: 0.5em 0; text-indent: 0; }}
.toc {{ font-size: 10pt; }}
.watermark {{ position: fixed; top: 40%; left: 10%; transform: rotate(-25deg);
  font-size: 28pt; color: rgba(180,60,40,0.12); pointer-events: none; z-index: 999; }}
</style>
</head>
<body>
<div class="watermark">MVP V2</div>
<div class="cover">
  <p class="badge">{badge}</p>
  <h1>{title}</h1>
  <p>{vol}</p>
  <p>Unit 1 · A001–A005 · Hybrid Voice V2.0</p>
  <p style="font-size:9pt;color:#666;">© 2026-06-08 · RGB试读 · 非印厂文件</p>
</div>
<div class="toc">
  <h2>{'目次' if is_ja else '目录'}</h2>
  <ul>"""
    ]
    for c in CASES:
        t = c["jp_title"] if is_ja else c["cn_title"]
        parts.append(f"<li>{c['id']} · {t}</li>")
    parts.append("</ul></div>")

    for c in CASES:
        fname = c["jp_file"] if is_ja else c["cn_file"]
        raw = strip_v2_reader(bde.read_text(BODY / fname))
        parts.append(f'<div class="case-block" id="{c["id"]}">')
        for stem, cap in c["illus"]:
            p = DEPTH / stem
            if not p.exists():
                p = first_glob(DEPTH, stem.replace(".png", "*"))
            parts.append(figure_block(p, cap))
        parts.append(prose_to_html(raw, lang))
        parts.append("</div>")

    parts.append("</body></html>")
    return "\n".join(parts)


def first_glob(folder: Path, pattern: str) -> Path | None:
    hits = sorted(folder.glob(pattern))
    return hits[0] if hits else None


def html_to_pdf(html_path: Path, pdf_path: Path) -> bool:
    return bde.html_to_pdf_chrome(html_path, pdf_path)


def copy_bodies_and_illus() -> list[str]:
    paths: list[str] = []
    for sub, label in [("01_五案CN", "cn"), ("02_五案JP", "jp")]:
        d = OUT_ROOT / sub
        d.mkdir(parents=True, exist_ok=True)
        for c in CASES:
            src = BODY / (c["cn_file"] if label == "cn" else c["jp_file"])
            dst = d / src.name
            shutil.copy2(src, dst)
            paths.append(str(dst.resolve()))

    ill = OUT_ROOT / "03_插图"
    ill.mkdir(parents=True, exist_ok=True)
    for c in CASES:
        sub = c.get("illus_subdir")
        case_dir = ill / sub if sub else ill
        case_dir.mkdir(parents=True, exist_ok=True)
        for stem, cap in c["illus"]:
            src = DEPTH / stem
            if src.exists():
                dst_flat = ill / src.name
                dst_case = case_dir / src.name
                shutil.copy2(src, dst_flat)
                if dst_case != dst_flat:
                    shutil.copy2(src, dst_case)
                paths.append(str(dst_flat.resolve()))
                if dst_case != dst_flat:
                    paths.append(str(dst_case.resolve()))
    return paths


def build_ppt() -> Path:
    try:
        from pptx import Presentation
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt
    except ImportError:
        subprocess.run([sys.executable, "-m", "pip", "install", "python-pptx", "-q"], check=True)
        from pptx import Presentation
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    blank = prs.slide_layouts[6]
    MX, MY = Inches(0.65), Inches(0.45)
    CW = Inches(14.7)

    def slide_title(title: str, sub: str = "") -> None:
        s = prs.slides.add_slide(blank)
        box = s.shapes.add_textbox(MX, Inches(2.2), CW, Inches(4))
        tf = box.text_frame
        tf.text = title
        tf.paragraphs[0].font.size = Pt(36)
        tf.paragraphs[0].font.bold = True
        if sub:
            p = tf.add_paragraph()
            p.text = sub
            p.font.size = Pt(18)

    def slide_bullets(title: str, lines: list[str]) -> None:
        s = prs.slides.add_slide(blank)
        tb = s.shapes.add_textbox(MX, MY, CW, Inches(0.8))
        tb.text_frame.text = title
        tb.text_frame.paragraphs[0].font.size = Pt(28)
        tb.text_frame.paragraphs[0].font.bold = True
        body = s.shapes.add_textbox(MX, Inches(1.2), CW, Inches(7))
        tf = body.text_frame
        tf.word_wrap = True
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.size = Pt(16)

    slide_title(
        "学堂趣事录 · Vol1 Unit1",
        "MVP V2.0 · A001–A005 · 2026-06-08 deadline delivery",
    )

    for c in CASES:
        slide_bullets(
            f"{c['id']} · {c['cn_title']}",
            [
                f"JP: {c['jp_title']}",
                c["hook"],
                "CN: REVIEW_LOOP R3 · Hybrid Voice V2.0",
                "Illus: G1draft / V2 depth anchor · DEMO removed 2026-06-08",
                "Status: MVP test sample · 非 G-IMG PRODUCT lock",
            ],
        )

    slide_bullets(
        "MVP 交付摘要",
        [
            "CN正文：五案 R3 reader-flow 完成",
            "JP：五案 G-JP MVP（A001 full · A002–A005 batch 2026-06-08）",
            "插图：MVP 10 关键帧（1 G1 c01 + 9 G1 PH）· CLASS_5-2 座席同步 · 案01–05 子目录",
            "PDF：中日各一 · reader.html 可浏览器打印",
            "Gap：G-JP LOCK · G-IMG PRODUCT · E06 PASS 待后续",
        ],
    )

    out = OUT_ROOT / PPT_NAME
    prs.save(str(out))
    return out


def write_manifest(extra_paths: list[str]) -> Path:
    manifest = OUT_ROOT / "00_MVP交付清单.md"
    lines = [
        "# MVP V2 Unit1 · 交付清单 · 2026-06-08",
        "",
        "> **状态**：MVP v3 · G1 扩展 + R5 CN · CLASS_5-2 座席同步 2026-06-08 · 非 G-IMG PRODUCT lock",
        "> **清理 SSOT**：`V2迁移/48_教室座位插图清理与更新_V0.1.md`",
        "> **进度 SSOT**：`V2迁移/32_MVP续作2_插图与R5_20260608.md`",
        "",
        "## 核心交付物",
        "",
        f"| 类型 | 绝对路径 |",
        f"|------|----------|",
        f"| PDF（中文·FULL读者版） | `{OUT_ROOT / PDF_FULL_CN_NAME}` |",
        f"| PDF（中文） | `{OUT_ROOT / PDF_NAME.replace('学堂趣事录', '学堂趣事录_CN').replace('.pdf', '')}` |",
        f"| PDF | `{OUT_ROOT / PDF_NAME}` |",
        f"| PPT | `{OUT_ROOT / PPT_NAME}` |",
        f"| HTML | `{OUT_ROOT / 'reader.html'}` |",
        f"| 阶段文档 | `{MIGRATION / '30_MVP截止交付_20260608.md'}` |",
        "",
        "## 正文源文件",
        "",
    ]
    for c in CASES:
        lines.append(f"- **{c['id']}** CN: `{BODY / c['cn_file']}`")
        lines.append(f"- **{c['id']}** JP: `{BODY / c['jp_file']}`")
    lines.extend(["", "## 包内副本", ""])
    for p in sorted(OUT_ROOT.rglob("*")):
        if p.is_file():
            lines.append(f"- `{p.resolve()}`")
    for p in extra_paths:
        if p not in lines:
            lines.append(f"- `{p}`")
    manifest.write_text("\n".join(lines), encoding="utf-8")
    return manifest


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--skip-pdf", action="store_true")
    args = parser.parse_args()

    OUT_ROOT.mkdir(parents=True, exist_ok=True)
    copied = copy_bodies_and_illus()

    # Combined reader: JP primary (market), CN in separate PDF
    html_ja = OUT_ROOT / "reader_ja.html"
    html_cn = OUT_ROOT / "reader_cn.html"
    html_ja.write_text(build_reader_html("ja"), encoding="utf-8")
    html_cn.write_text(build_reader_html("zh"), encoding="utf-8")
    reader = OUT_ROOT / "reader.html"
    shutil.copy2(html_ja, reader)

    pdf_path = OUT_ROOT / PDF_NAME
    pdf_cn = OUT_ROOT / PDF_NAME.replace("学堂趣事录", "学堂趣事录_CN")
    pdf_full_cn = OUT_ROOT / PDF_FULL_CN_NAME
    if not args.skip_pdf:
        if not html_to_pdf(html_ja, pdf_path):
            print("JP PDF failed — open reader.html", file=sys.stderr)
        if not html_to_pdf(html_cn, pdf_cn):
            print("CN PDF failed", file=sys.stderr)
        if pdf_cn.exists() and not html_to_pdf(html_cn, pdf_full_cn):
            # fallback: copy CN MVP PDF as FULL reader edition
            shutil.copy2(pdf_cn, pdf_full_cn)
            print(f"FULL CN PDF copied from MVP CN: {pdf_full_cn}")

    ppt_path = build_ppt()
    manifest = write_manifest(copied + [str(ppt_path), str(pdf_path)])

    print(f"OK: {manifest}")
    print(f"PDF: {pdf_path}")
    print(f"PPT: {ppt_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
