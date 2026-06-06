#!/usr/bin/env python3
"""V0.1-A · 第一话审阅包 · 序 + A001 · 日本語主 + 中文附录.

Outputs (05_出版成果/):
  - 学堂奇事録_Vol1_V0.1-A_第一話审阅包_日本語_{date}.pdf
  - 学堂奇事録_Vol1_V0.1-A_出版提案_{date}.pptx
  - V01A_第一話审阅包_门禁对照表_{date}.md
"""

from __future__ import annotations

import re
import subprocess
import sys
from datetime import date
from pathlib import Path

TOOLS = Path(__file__).resolve().parent
FORMAL = TOOLS.parents[1]
SAMPLE = FORMAL.parent / "样章包"
OUT = FORMAL / "05_出版成果"
DISPLAY_TOOLS = FORMAL / "04_展示版" / "tools"

JP_BODY = SAMPLE / "04_样章_序+案01_正文_日本語.txt"
CN_BODY = FORMAL / "01_正文" / "案01_翘边的海报_HybridVoice_V1.1定稿.txt"
HIKARU_TXT = SAMPLE / "05_陸瑆日记页_样章.txt"
EXPERIMENT = SAMPLE / "12_案01_家庭实验页_样章.txt"
CLUE_MD = SAMPLE / "07_线索卡_设计稿.md"
CHARS = SAMPLE / "08_角色介绍_一页.md"

DEPTH = SAMPLE / "插图" / "depth_anchor"
SUPP = DEPTH / "supplement"
ILLUST = SAMPLE / "插图"

OUT_DATE = date.today().strftime("%Y%m%d")
PDF_NAME = f"学堂奇事録_Vol1_V0.1-A_第一話审阅包_日本語_{OUT_DATE}.pdf"
HTML_NAME = f"学堂奇事録_Vol1_V0.1-A_第一話审阅包_日本語_{OUT_DATE}.html"
PPT_NAME = f"学堂奇事録_Vol1_V0.1-A_出版提案_{OUT_DATE}.pptx"
GATE_MD = f"V01A_第一話审阅包_门禁对照表_{OUT_DATE}.md"

sys.path.insert(0, str(DISPLAY_TOOLS))
import build_display_edition as bde  # noqa: E402

# (chapter keyword, section num, [(path relative candidates, css, caption ja, status)])
SHOT_INJECT: list[tuple[str, str | None, list[tuple[list[Path], str, str, str]]]] = [
    (
        "序",
        None,
        [
            (
                [SUPP / "V-S00_序钩_海报微翘_v0.1.png"],
                "half",
                "序 · ポスターの端（V-S00 · 探索 v0.1）",
                "🟢",
            )
        ],
    ),
    (
        "めくれたポスター",
        "1",
        [
            (
                [DEPTH / "V-S01-A1_侧廊发现_v0.2.png", DEPTH / "V-S01-A1_侧廊发现.png"],
                "half",
                "DA1 · 側廊で発見（L0 · v0.2）",
                "🟢",
            )
        ],
    ),
    (
        "めくれたポスター",
        "2",
        [
            (
                [DEPTH / "V-S01-A2_误导搜查_v0.2.png", DEPTH / "V-S01-A2_误导搜查.png"],
                "half",
                "DA2 · 誤った捜査（L0 · v0.2）",
                "🟢",
            ),
            (
                [SUPP / "V-S01-A2b_珣指风侧_MCU_v0.1.png"],
                "quarter",
                "DA2b · 珣「風の側」（inset · v0.1）",
                "🟢",
            ),
        ],
    ),
    (
        "めくれたポスター",
        "3",
        [
            (
                [DEPTH / "V-S01-A3_风侧线索_v0.2.png", DEPTH / "V-S01-A3_风侧线索.png"],
                "half",
                "DA3 · 風側の手がかり（L0 · v0.2）",
                "🟢",
            ),
            (
                [SUPP / "V-S01-A3b_胶带开边_ECUs_v0.1.png"],
                "quarter",
                "DA3b · テープの開き辺（ECU · v0.1）",
                "🟢",
            ),
        ],
    ),
    (
        "めくれたポスター",
        "4",
        [
            (
                [DEPTH / "V-S01-A4_验证收束_v0.2.png", DEPTH / "V-S01-A4_验证收束.png"],
                "half",
                "DA4 · 検証と収束（L0 · v0.2）",
                "🟢",
            ),
            (
                [SUPP / "V-S01-A4b_换贴对比_MCU_v0.1.png"],
                "quarter",
                "DA4b · 貼り方の比較（MCU · v0.1）",
                "🟢",
            ),
            (
                [DEPTH / "V-S01-TAIL_壁报空栏_v0.2.png", DEPTH / "V-S01-TAIL_壁报空栏.png"],
                "quarter",
                "DC1 · 壁報の空欄（L0 · v0.2）",
                "🟢",
            ),
        ],
    ),
]


def first_existing(candidates: list[Path]) -> Path | None:
    for p in candidates:
        if p.exists():
            return p
    return None


def figure_from_paths(candidates: list[Path], css: str, caption: str, status: str) -> str:
    path = first_existing(candidates)
    badge = f'<span class="shot-badge">{status} 探索稿·非印厂</span>'
    if path:
        uri = bde.img_data_uri(path)
        return (
            f'<figure class="illus {css}"><img src="{uri}" alt="{caption}"/>'
            f'<figcaption>{caption} {badge}</figcaption></figure>'
        )
    names = " / ".join(p.name for p in candidates)
    return (
        f'<div class="placeholder-illus {css}"><span>⬜ 挿画待ち · {names}</span>'
        f'<p class="cap">{caption}</p></div>'
    )


def inject_a001_shots(chapter: str, section: str | None) -> str:
    sec_num = None
    if section and section[0].isdigit():
        sec_num = section.split()[0]
    parts: list[str] = []
    for key, rule_sec, figs in SHOT_INJECT:
        if key not in chapter:
            continue
        if rule_sec is None and section is not None:
            continue
        if rule_sec is not None and sec_num != rule_sec:
            continue
        for paths, css, cap, st in figs:
            parts.append(figure_from_paths(paths, css, cap, st))
    return "\n".join(parts)


def prose_to_html_a001(text: str) -> str:
    blocks: list[str] = []
    para: list[str] = []
    current_chapter = ""

    for line in text.splitlines():
        s = line.rstrip()
        if not s:
            if para:
                blocks.append("<p>" + "<br/>".join(para) + "</p>")
                para = []
            continue
        if s == "---":
            blocks.append('<hr class="scene-break"/>')
            continue
        if s.startswith("=") and len(s.strip()) > 10:
            continue
        if re.match(r"^(序|一、|めくれた|翘边)", s):
            current_chapter = s.split("-")[0].strip()
            blocks.append(f'<h2 class="chapter">{current_chapter}</h2>')
            fig = inject_a001_shots(current_chapter, None)
            if fig:
                blocks.append(fig)
            continue
        if s.startswith("### "):
            sec = s[4:]
            blocks.append(f'<h3 class="section">{sec}</h3>')
            fig = inject_a001_shots(current_chapter, sec)
            if fig:
                blocks.append(fig)
            continue
        if s.startswith("  ") and not s.startswith("  http"):
            para.append(f'<span class="indent">{s.strip()}</span>')
        elif re.match(r"^-{4,}$", s):
            continue
        else:
            clean = re.sub(r"\*\*(.+?)\*\*", r"<strong>\1</strong>", s)
            if clean.strip() and not clean.startswith("【"):
                para.append(clean)
    if para:
        blocks.append("<p>" + "<br/>".join(para) + "</p>")
    return "\n".join(blocks)


def gate_rows() -> list[tuple[str, str, str, str]]:
    return [
        ("正典", "第一話=A001 非湿椅子", "✅ LOCK", "序+めくれたポスター"),
        ("正典", "Vol1=Plan B", "✅ LOCK", "本包は事件①のみ"),
        ("Gate A", "P0-01/02 IP LOCK", "✅", "2026-06-05"),
        ("C03", "公平线索≥3", "✅", "風側/雨天/貼り方"),
        ("C04", "温柔真相", "✅", "志郎非悪意"),
        ("V13", "桥梁书年龄感", "🟡", "v0.1/v0.2 探索稿"),
        ("E04", "日文文体", "🟡", "JP_VOICE_v1 · 田中待ち"),
        ("E07", "校园五维", "⬜", "側廊/クラブ/壁報"),
        ("科学", "P0-04 A001", "⬜", "DB1+DA3 待签"),
        ("画风", "正篇Sheet", "🟡", "六帧+增补 v0.1/v0.2"),
        ("S20", "非湿椅子叙事", "✅", "实验=テープ+ドライヤー"),
    ]


def gate_table_html() -> str:
    rows = "".join(
        f"<tr><td>{a}</td><td>{b}</td><td class='st'>{c}</td><td>{d}</td></tr>"
        for a, b, c, d in gate_rows()
    )
    return f"""
    <div class="page">
      <h2>正典 · V0.1-A 门禁对照</h2>
      <p class="lead">第一話里程碑 · 序 + A001 · 专家/编辑/美工并行审稿用</p>
      <table class="gate-table">
        <thead><tr><th>項目</th><th>区分</th><th>状態</th><th>備考</th></tr></thead>
        <tbody>{rows}</tbody>
      </table>
    </div>
    """


def shot_map_html() -> str:
    rows = [
        ("V-S00", "SC-00", "序钩", "P2", "MS", "✅ v0.1"),
        ("DA1", "SC-01-1", "发现", "P0", "MS", "✅ v0.2"),
        ("DA2", "SC-01-2", "误导", "P1", "MS", "✅ v0.2"),
        ("DA2b", "SC-01-2", "珣指风", "P1", "MCU", "✅ v0.1"),
        ("DA3", "SC-01-3", "线索", "P0", "CU", "✅ v0.2"),
        ("DA3b", "SC-01-3", "开边ECU", "P0", "ECU", "✅ v0.1"),
        ("DA4", "SC-01-4", "收束", "P0", "MS", "✅ v0.2"),
        ("DA4b", "SC-01-4", "换贴", "P0", "MCU", "✅ v0.1"),
        ("DC1", "SC-01-5", "尾钩", "P1", "CU", "✅ v0.2"),
        ("DB1", "—", "机制SUM", "P0", "信息图", "✅ v0.2"),
        ("V-S02", "案末", "瑆日记", "P1", "B轨", "✅ v0.1"),
    ]
    trs = "".join(
        f"<tr><td>{a}</td><td>{b}</td><td>{c}</td><td>{d}</td><td>{e}</td><td>{f}</td></tr>"
        for a, b, c, d, e, f in rows
    )
    return f"""
    <div class="page">
      <h2>分镜 Shot Map · A001（V1.2）</h2>
      <p class="lead">詳細：<code>样章包/03_案01_分镜头与插页地图_V1.2_完整分镜.md</code></p>
      <table class="gate-table shot-map">
        <thead><tr><th>Shot</th><th>SC</th><th>機能</th><th>P</th><th>景別</th><th>状態</th></tr></thead>
        <tbody>{trs}</tbody>
      </table>
    </div>
    """


def review_columns_html() -> str:
    roles = [
        ("文字编辑 E04", "注音·文体·年级适配"),
        ("科学顾问", "DB1/DA3/公平线索/实验安全"),
        ("插画·美工 V13", "画风Sheet·角色一致·占位"),
        ("日本文化 E07", "侧廊·上履き·クラブ·壁報"),
        ("主编", "Gate A PASS / 进 V0.2"),
    ]
    blocks = ""
    for title, hint in roles:
        blocks += f"""
        <div class="review-block">
          <h3>{title}</h3>
          <p class="hint">{hint}</p>
          <div class="review-lines">
            <p>□ PASS &nbsp; □ HOLD &nbsp; □ 修正要</p>
            <p class="line">意见：_____________________________________________</p>
            <p class="line">署名：______________ 日期：______________</p>
          </div>
        </div>
        """
    return f'<div class="page review-page"><h2>并行审核栏 · V0.1-A</h2>{blocks}</div>'


def hikaru_page_html() -> str:
    img = first_existing(
        [
            ILLUST / "V-S02_瑆日记页_v0.1.png",
            ILLUST / "V-S02_陸瑆日记页.png",
        ]
    )
    img_html = ""
    if img:
        img_html = figure_from_paths([img], "half", "V-S02 · 瑆のノート（B轨 · v0.1）", "🟢")
    body = bde.strip_txt_meta(bde.read_text(HIKARU_TXT)) if HIKARU_TXT.exists() else ""
    text_html = prose_to_html_a001(body) if body else ""
    return f"""
    <div class="page diary-page">
      <h2>陸瑆（ひかる）のノート · 样章页</h2>
      {img_html}
      {text_html}
      <p class="note">B轨：正文层と画风を明確に区別 · 问句收束 · 不剧透</p>
    </div>
    """


def db1_and_experiment_html() -> str:
    db1 = figure_from_paths(
        [DEPTH / "V-S01-B1_风侧机制图_v0.2.png", DEPTH / "V-S01-B1_风侧机制图.png"],
        "half",
        "DB1 · 風側メカニズム SUM（P0 · v0.2）",
        "🟢",
    )
    exp = bde.read_text(EXPERIMENT) if EXPERIMENT.exists() else ""
    exp_html = prose_to_html_a001(re.sub(r"^# .+\n", "", exp))
    return f"""
    <div class="page">
      <h2>家庭実験 · 事件①</h2>
      <p class="lead">カード紙+テープ+ドライヤー · 大人同伴 · 非湿椅子实验</p>
      {db1}
      <div class="experiment-text">{exp_html}</div>
    </div>
    """


def clue_card_html() -> str:
    raw = bde.read_text(CLUE_MD) if CLUE_MD.exists() else ""
    pre = re.search(r"```\n([\s\S]*?)```", raw)
    box = pre.group(1) if pre else "线索卡 #01 · 见 07_线索卡_设计稿.md"
    return f"""
    <div class="page clue-card">
      <h2>観察クラブ · 手がかりカード #01</h2>
      <pre class="clue-box">{box.strip()}</pre>
      <p class="note">三项线索正文均已出现 · 仅回顾 · ⬜ 版式 E22 待决</p>
    </div>
    """


def cn_appendix_html() -> str:
    raw = bde.strip_txt_meta(bde.read_text(CN_BODY)) if CN_BODY.exists() else ""
    body = prose_to_html_a001(raw)
    return f"""
    <div class="page appendix-cn">
      <h2>中文附录 · 对照用（非主交付语言）</h2>
      <p class="lead">源：<code>01_正文/案01_翘边的海报_HybridVoice_V1.1定稿.txt</code> · V1.1</p>
      <div class="prose-cn">{body}</div>
    </div>
    """


def version_page_html() -> str:
    return f"""
    <div class="page">
      <h2>版本记录 · V0.1-A</h2>
      <table class="gate-table">
        <tr><th>版</th><td>V0.1-A · 第一話审阅包</td></tr>
        <tr><th>日期</th><td>{OUT_DATE}</td></tr>
        <tr><th>范围</th><td>序 + A001 · 非全卷 · 非湿椅子</td></tr>
        <tr><th>主语言</th><td>日本語 + 中文附录</td></tr>
        <tr><th>插图</th><td>L0 v0.2 + supplement v0.1 · 探索稿</td></tr>
        <tr><th>下一版</th><td>V0.2：专家意见汇总 · L0 v1.0 精修 · P0-04 PASS</td></tr>
      </table>
      <h3>待修改项（模板）</h3>
      <ul>
        <li>⬜ E07 田中校园督查</li>
        <li>⬜ P0-04 科学顾问签 DB1</li>
        <li>⬜ 画风 Sheet 六帧主编 PASS</li>
        <li>⬜ 日文 E04 全文推敲</li>
      </ul>
      <p class="note">rebuild: python 05_出版成果/tools/build_v01a_episode1_review_pack.py</p>
    </div>
    """


def build_html() -> str:
    jp_raw = bde.strip_txt_meta(bde.read_text(JP_BODY))
    body = prose_to_html_a001(jp_raw)
    cover_img = first_existing(
        [DEPTH / "V-S01-A1_侧廊发现_v0.2.png", ILLUST / "V-S01_侧廊海报.png"]
    )
    cover_uri = bde.img_data_uri(cover_img) if cover_img else ""

    chars = (
        "【事件① 出场】\n"
        "陸珣（しゅん）· 5年2組 · 転校生\n"
        "伊藤光 · 5年2組 · おもしろ観察クラブ\n"
        "加藤慧美 · 5年1組 · 記録\n"
        "松本志郎 · 5年3組 · 検証\n"
        "陸瑆（ひかる）· 4年2組 · 日記层（B轨）\n"
        "※ 理紗/中谷/父母：本话不出场"
    )

    return f"""<!DOCTYPE html>
<html lang="ja">
<head>
<meta charset="utf-8"/>
<title>学堂奇事録 V0.1-A 第一話审阅包</title>
<style>
@page {{ size: A5; margin: 14mm 11mm 18mm 14mm; }}
* {{ box-sizing: border-box; }}
body {{
  font-family: "Yu Gothic", "Meiryo", "Hiragino Sans", sans-serif;
  font-size: 10.5pt; line-height: 1.58; color: #1a1a1a; max-width: 148mm; margin: 0 auto;
}}
.page {{ page-break-after: always; padding: 3mm 0; min-height: 210mm; position: relative; }}
.page:last-child {{ page-break-after: auto; }}
.review-margin {{
  position: fixed; right: 0; top: 0; width: 22mm; height: 100%;
  border-left: 1px dotted #ccc; font-size: 7pt; color: #aaa;
  writing-mode: vertical-rl; padding: 4mm 2mm;
}}
.cover {{
  display:flex; flex-direction:column; justify-content:center; align-items:center;
  text-align:center; min-height:240mm;
  background: linear-gradient(165deg, #faf7f2 0%, #e8e0d4 100%);
}}
.cover h1 {{ font-size: 22pt; margin: 0.3em 0; }}
.cover .vol {{ font-size: 12pt; color: #444; }}
.cover .badge {{
  margin-top: 1em; font-size: 9pt; border: 2px solid #8a6d4b;
  padding: 0.4em 0.9em; color: #5a4030; background: #fffef8;
}}
.cover img {{ max-width: 68%; margin: 0.8em 0; border: 1px solid #ccc; }}
h2 {{ font-size: 12pt; border-bottom: 1px solid #bbb; padding-bottom: 0.2em; }}
h3.section {{ font-size: 10pt; color: #555; }}
.lead, .note {{ font-size: 8.5pt; color: #666; }}
.gate-table {{ width:100%; border-collapse: collapse; font-size: 8.5pt; margin: 0.6em 0; }}
.gate-table th, .gate-table td {{ border: 1px solid #ccc; padding: 0.3em 0.4em; vertical-align: top; }}
.gate-table .st {{ font-weight: bold; white-space: nowrap; }}
.illus {{ text-align:center; margin: 0.5em 0; page-break-inside: avoid; }}
.illus img {{ max-width: 100%; height: auto; }}
.illus.half img {{ max-height: 78mm; }}
.illus.quarter img {{ max-height: 46mm; }}
.illus figcaption {{ font-size: 7.5pt; color: #555; }}
.shot-badge {{ color: #8a6d4b; font-size: 7pt; }}
.placeholder-illus {{
  border: 2px dashed #bbb; min-height: 40mm; display:flex; flex-direction:column;
  align-items:center; justify-content:center; color:#999; font-size:9pt; margin:0.5em 0;
}}
.diary-page {{
  background: #fffef8; padding: 1em; border-left: 3px solid #d4c4a8;
  background-image: repeating-linear-gradient(#0000 0 1.3em, #e8e4dc 1.3em 1.31em);
}}
.review-block {{ margin: 0.8em 0; padding: 0.5em; border: 1px solid #ddd; }}
.review-block .hint {{ font-size: 8pt; color: #777; }}
.review-block .line {{ border-bottom: 1px dotted #ccc; min-height: 1.4em; }}
.clue-box {{
  font-family: "Yu Gothic", monospace; font-size: 9pt;
  background: #f7f4ef; padding: 1em; border: 1px solid #ccc; white-space: pre-wrap;
}}
.appendix-cn {{ background: #fafafa; }}
.appendix-cn .prose-cn {{ font-family: "Microsoft YaHei", sans-serif; }}
p {{ margin: 0.45em 0; }}
p .indent {{ display:block; padding-left:1.2em; }}
hr.scene-break {{ border:none; border-top:1px dotted #ccc; margin:0.8em 0; }}
@media print {{ body {{ max-width: none; }} }}
</style>
</head>
<body>

<div class="page cover">
  <p class="badge">V0.1-A · 第一話审阅包 · 非売品 · RGB探索稿</p>
  <h1>学堂奇事録</h1>
  <p class="vol">第1巻 · 事件①「めくれたポスター」<br/>序 + A001 · Plan B</p>
  {"<img src='" + cover_uri + "' alt='DA1'/>" if cover_uri else ""}
  <p class="note">{OUT_DATE} · 日本語主 · 中文附录 · Gate A</p>
</div>

<div class="page">
  <h2>出版声明</h2>
  <ul style="font-size:9.5pt;line-height:1.55;">
    <li><strong>種別</strong>：V0.1-A 团队共审里程碑（专家库/编辑/美工）</li>
    <li><strong>非</strong>：印厂原稿 · CMYK终稿 · 湿椅子(C001)卷</li>
    <li><strong>正典</strong>：Vol1 方案B · 第一話=A001 · IP LOCK 2026-06-05</li>
  </ul>
</div>

{gate_table_html()}
{shot_map_html()}

<div class="page">
  <h2>登場人物（事件①）</h2>
  <pre style="white-space:pre-wrap;font-family:inherit;font-size:10pt;">{chars}</pre>
</div>

<div class="page">
  <h2>目次 · V0.1-A</h2>
  <ul class="toc" style="list-style:none;padding:0;">
    <li>序 · 四月の第二月曜日</li>
    <li>一、めくれたポスター（§1–4 + 挿画）</li>
    <li>陸瑆のノート · V-S02</li>
    <li>手がかりカード #01</li>
    <li>DB1 + 家庭実験</li>
    <li>并行审核栏</li>
    <li>中文附录</li>
    <li>版本记录</li>
  </ul>
</div>

<div class="page prose">
  {body}
</div>

{hikaru_page_html()}
{clue_card_html()}
{db1_and_experiment_html()}
{review_columns_html()}
{cn_appendix_html()}
{version_page_html()}

</body>
</html>
"""


def write_gate_md() -> Path:
    lines = [
        f"# V0.1-A 第一話审阅包 · 门禁对照 · {OUT_DATE}",
        "",
        "| 项 | 类别 | 状态 | 备注 |",
        "|----|------|------|------|",
    ]
    for row in gate_rows():
        lines.append(f"| {row[0]} | {row[1]} | {row[2]} | {row[3]} |")
    lines.extend(["", "生成：`build_v01a_episode1_review_pack.py`", ""])
    path = OUT / GATE_MD
    path.write_text("\n".join(lines), encoding="utf-8")
    return path


def build_ppt(pdf_hint: Path) -> Path | None:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        subprocess.run([sys.executable, "-m", "pip", "install", "python-pptx", "-q"], check=True)
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    blank = prs.slide_layouts[6]
    MX = Inches(0.65)

    def title(t: str, sub: str = "") -> None:
        s = prs.slides.add_slide(blank)
        box = s.shapes.add_textbox(MX, Inches(2.2), Inches(14.7), Inches(4))
        tf = box.text_frame
        tf.text = t
        tf.paragraphs[0].font.size = Pt(36)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        if sub:
            p = tf.add_paragraph()
            p.text = sub
            p.font.size = Pt(18)
            p.alignment = PP_ALIGN.CENTER

    def bullets(t: str, lines: list[str]) -> None:
        s = prs.slides.add_slide(blank)
        tb = s.shapes.add_textbox(MX, Inches(0.5), Inches(14.7), Inches(8))
        tf = tb.text_frame
        tf.text = t
        tf.paragraphs[0].font.size = Pt(28)
        tf.paragraphs[0].font.bold = True
        for line in lines:
            p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(16)
            p.space_before = Pt(8)

    title("V0.1-A · 第一話审阅包", "序 + A001 めくれたポスター · 日本語主")
    bullets("本包定位", [
        "• Gate A 里程碑 · 非全卷 · 非湿椅子",
        "• 专家/编辑/美工并行审稿基准",
        "• 插图=探索稿 v0.1/v0.2 · RGB",
        f"• PDF：{pdf_hint.name}",
    ])
    bullets("Shot Map 要点", [
        "L0: DA1–DA4 + DB1 + DC1（6帧）",
        "L1: V-S02 瑆 + DA2b/DA3b/DA4b",
        "L2: V-S00 序钩",
        "待签: P0-04 科学 · E07 田中 · 画风Sheet",
    ])
    img = first_existing([DEPTH / "V-S01-A1_侧廊发现_v0.2.png"])
    if img:
        s = prs.slides.add_slide(blank)
        s.shapes.add_textbox(MX, Inches(0.4), Inches(6), Inches(0.5)).text_frame.text = "DA1 · v0.2"
        s.shapes.add_picture(str(img), Inches(5.5), Inches(0.6), width=Inches(10))

    bullets("审核分工", [
        "文字 E04 · 科学 · 插画 V13 · 文化 E07 · 主编",
        "PDF 内「并行审核栏」填写",
        "汇总 → V0.2 迭代",
    ])

    out = OUT / PPT_NAME
    prs.save(str(out))
    return out


def main() -> int:
    if not JP_BODY.exists():
        print(f"Missing: {JP_BODY}", file=sys.stderr)
        return 1
    OUT.mkdir(parents=True, exist_ok=True)

    html = build_html()
    html_path = OUT / HTML_NAME
    html_path.write_text(html, encoding="utf-8")
    print(f"HTML: {html_path}")

    pdf_path = OUT / PDF_NAME
    if not bde.html_to_pdf_chrome(html_path, pdf_path):
        print("PDF failed — open HTML and print", file=sys.stderr)
        return 1
    print(f"PDF: {pdf_path} ({pdf_path.stat().st_size // 1024} KB)")

    md_path = write_gate_md()
    print(f"Gate: {md_path}")

    ppt = build_ppt(pdf_path)
    if ppt:
        print(f"PPT: {ppt} ({ppt.stat().st_size // 1024} KB)")
    return 0


if __name__ == "__main__":
    sys.exit(main())
