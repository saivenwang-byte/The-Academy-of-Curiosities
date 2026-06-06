#!/usr/bin/env python3
"""Vol1 expert publication pack: V0.1-A (序+A001) + V0.1-B (全卷提案).

Modes (--mode):
  v01a  — 第一話审阅包 · JP main + CN appendix
  v01b  — Vol1 全卷提案包 · 5案骨架 · 出版社对外
  both  — default · 生成两套 PDF + PPT + gate MD

Outputs (05_出版成果/):
  V0.1-A:
    学堂奇事録_Vol1_V0.1-A_第一話审阅包_日本語_{date}.pdf
    学堂奇事録_Vol1_V0.1-A_出版提案_{date}.pptx
    V01A_第一話审阅包_门禁对照表_{date}.md
  V0.1-B:
    学堂奇事録_Vol1_V0.1-B_全卷提案包_日本語_{date}.pdf
    学堂奇事録_Vol1_V0.1-B_出版提案_{date}.pptx
    V01B_全卷提案包_门禁对照表_{date}.md
"""

from __future__ import annotations

import argparse
import re
import subprocess
import sys
from datetime import date
from pathlib import Path

TOOLS = Path(__file__).resolve().parent
FORMAL = TOOLS.parents[1]
SAMPLE = FORMAL.parent / "样章包"
OUT = FORMAL / "05_出版成果"
DISPLAY_TOOLS = FORMAL / "04_展示版" / "tools"
ILLUST_TOOLS = FORMAL / "02_插画" / "tools"
BODY_DIR = FORMAL / "01_正文"
CANON_JP = FORMAL / "01_正本" / "学堂趣事录_Vol1_觉得奇怪就先观察_正本_日本語.txt"

JP_A001 = BODY_DIR / "案01_めくれたポスター_V1.1定稿_日本語.txt"
CN_A001 = BODY_DIR / "案01_翘边的海报_HybridVoice_V1.1定稿.txt"
JP_SAMPLE = SAMPLE / "04_样章_序+案01_正文_日本語.txt"
SHOT_MAP_MD = SAMPLE / "03_案01_分镜头与插页地图_V1.2_完整分镜.md"
HIKARU_TXT = SAMPLE / "05_陸瑆日记页_样章.txt"
EXPERIMENT = SAMPLE / "12_案01_家庭实验页_样章.txt"
CLUE_MD = SAMPLE / "07_线索卡_设计稿.md"
NOTES_JP = FORMAL / "03_笔记" / "陸瑆笔记_Vol1_日本語.md"

DEPTH = SAMPLE / "插图" / "depth_anchor"
SUPP = DEPTH / "supplement"
ILLUST_SAMPLE = SAMPLE / "插图"
ILLUST = FORMAL / "02_插画" / "assets"

OUT_DATE = date.today().strftime("%Y%m%d")

V01A = {
    "pdf": f"学堂奇事録_Vol1_V0.1-A_第一話审阅包_日本語_{OUT_DATE}.pdf",
    "html": f"学堂奇事録_Vol1_V0.1-A_第一話审阅包_日本語_{OUT_DATE}.html",
    "ppt": f"学堂奇事録_Vol1_V0.1-A_出版提案_{OUT_DATE}.pptx",
    "gate": f"V01A_第一話审阅包_门禁对照表_{OUT_DATE}.md",
}
V01B = {
    "pdf": f"学堂奇事録_Vol1_V0.1-B_全卷提案包_日本語_{OUT_DATE}.pdf",
    "html": f"学堂奇事録_Vol1_V0.1-B_全卷提案包_日本語_{OUT_DATE}.html",
    "ppt": f"学堂奇事録_Vol1_V0.1-B_出版提案_{OUT_DATE}.pptx",
    "gate": f"V01B_全卷提案包_门禁对照表_{OUT_DATE}.md",
}

CASES = [
    ("A001", "①", "めくれたポスター", "湿度·気流·貼り方", "案01_めくれたポスター_V1.1定稿_日本語.txt", "✅"),
    ("A002", "②", "ズレた泥のあと", "動線·側門の土", "案02_ずれた泥のあと_V1.1定稿_日本語.txt", "✅"),
    ("A003", "③", "空いている欄", "確認待ち·掲示規則", "案03_空いているマス_V1.1定稿_日本語.txt", "✅"),
    ("A004", "④", "チョーク粉の輪", "静電·掃除順", "案04_チョークの粉の輪_V1.1定稿_日本語.txt", "🟡"),
    ("A005", "⑤", "消しゴムのカス", "摩擦·入社", "案05_消しゴムの屑の向き_V1.1定稿_日本語.txt", "🟡"),
]

# A001 shot injection: (chapter key, section num, [(path candidates, css, caption, shot_id)])
SHOT_INJECT: list[tuple[str, str | None, list[tuple[list[Path], str, str, str]]]] = [
    ("序", None, [([SUPP / "V-S00_序钩_海报微翘_v0.1.png"], "half", "V-S00 · 序钩 · ポスター微翘", "V-S00")]),
    ("めくれたポスター", "1", [([DEPTH / "V-S01-A1_侧廊发现_v0.2.png", DEPTH / "V-S01-A1_侧廊发现.png"], "half", "DA1 · 側廊で発見", "DA1")]),
    ("めくれたポスター", "2", [
        ([DEPTH / "V-S01-A2_误导搜查_v0.2.png", DEPTH / "V-S01-A2_误导搜查.png"], "half", "DA2 · 誤った捜査", "DA2"),
        ([SUPP / "V-S01-A2b_珣指风侧_MCU_v0.1.png"], "quarter", "DA2b · 珣「風の側」", "DA2b"),
    ]),
    ("めくれたポスター", "3", [
        ([DEPTH / "V-S01-A3_风侧线索_v0.2.png", DEPTH / "V-S01-A3_风侧线索.png"], "half", "DA3 · 風側の手がかり", "DA3"),
        ([SUPP / "V-S01-A3b_胶带开边_ECUs_v0.1.png"], "quarter", "DA3b · テープの開き辺", "DA3b"),
    ]),
    ("めくれたポスター", "4", [
        ([DEPTH / "V-S01-A4_验证收束_v0.2.png", DEPTH / "V-S01-A4_验证收束.png"], "half", "DA4 · 検証と収束", "DA4"),
        ([SUPP / "V-S01-A4b_换贴对比_MCU_v0.1.png"], "quarter", "DA4b · 貼り方の比較", "DA4b"),
        ([DEPTH / "V-S01-TAIL_壁报空栏_v0.2.png", DEPTH / "V-S01-TAIL_壁报空栏.png"], "quarter", "DC1 · 壁報の空欄", "DC1"),
    ]),
]

# V0.1-B per-case placeholder art (non-A001)
CASE_ART: dict[str, list[tuple[str, str]]] = {
    "A002": [("V-S02_陸瑆日记页.png", "瑆ノート · 風とポスター"), ("V-S02-TAIL_无署名窄条.png", "空欄の無署名投稿")],
    "A003": [("V-S03-TAIL_投稿背面路线.png", "投稿の裏 · ルート")],
    "A004": [("V-S04_讲台粉笔灰圈.png", "教壇 · 粉の輪"), ("V-S04-TAIL_中谷远景.png", "中谷の半句")],
    "A005": [("V-S05_橡皮屑方向.png", "美術室 · カスの向き"), ("V-S05-TAIL_瑆日记扶正书.png", "瑆 · 本の角")],
}

sys.path.insert(0, str(DISPLAY_TOOLS))
import build_display_edition as bde  # noqa: E402


def first_existing(candidates: list[Path]) -> Path | None:
    for p in candidates:
        if p.exists():
            return p
    return None


def resolve_depth_png(png_id: str) -> list[Path]:
    """Map Shot Map PNG ID to candidate paths."""
    base = png_id.strip("`").strip()
    if base.startswith("supplement/"):
        name = Path(base).name
        return [SUPP / name]
    if base.startswith("插图/"):
        name = Path(base).name
        return [ILLUST_SAMPLE / name, ILLUST / name.replace("_v0.1", "")]
    stem = base
    return [
        DEPTH / f"{stem}_v0.2.png",
        DEPTH / f"{stem}.png",
        DEPTH / stem,
        ILLUST_SAMPLE / f"{stem}.png",
    ]


def figure_from_paths(
    candidates: list[Path],
    css: str,
    caption: str,
    shot_id: str = "",
    exploration: bool = True,
) -> str:
    path = first_existing(candidates)
    badge = '<span class="shot-badge">探索稿·非印刷</span>' if exploration else ""
    if path:
        uri = bde.img_data_uri(path)
        sid = f" [{shot_id}]" if shot_id else ""
        return (
            f'<figure class="illus {css}"><img src="{uri}" alt="{caption}"/>'
            f'<figcaption>{caption}{sid} {badge}</figcaption></figure>'
        )
    names = " / ".join(p.name for p in candidates)
    brief = caption or shot_id or names
    return (
        f'<div class="placeholder-illus {css}">'
        f'<span class="ph-id">{shot_id or "⬜"}</span>'
        f'<span class="ph-brief">{brief}</span>'
        f'<p class="cap">挿画待ち · {names}</p></div>'
    )


def placeholder_box(shot_id: str, brief: str, css: str = "quarter") -> str:
    return (
        f'<div class="placeholder-illus {css}">'
        f'<span class="ph-id">{shot_id}</span>'
        f'<span class="ph-brief">{brief}</span>'
        f'<p class="cap">探索稿プレースホルダ</p></div>'
    )


def parse_shot_map_v12() -> list[dict[str, str]]:
    """Parse V1.2 full shot map table from markdown."""
    if not SHOT_MAP_MD.exists():
        return []
    raw = SHOT_MAP_MD.read_text(encoding="utf-8")
    rows: list[dict[str, str]] = []
    in_table = False
    headers: list[str] = []
    for line in raw.splitlines():
        if line.startswith("| Shot |"):
            in_table = True
            headers = [c.strip().strip("*") for c in line.split("|")[1:-1]]
            continue
        if not in_table:
            continue
        if not line.startswith("|"):
            if rows:
                break
            continue
        if re.match(r"^\|[-\s|]+\|$", line):
            continue
        cells = [c.strip().strip("*") for c in line.split("|")[1:-1]]
        if len(cells) < len(headers):
            cells += [""] * (len(headers) - len(cells))
        rows.append(dict(zip(headers, cells)))
    return rows


def shot_map_table_html(compact: bool = False) -> str:
    rows = parse_shot_map_v12()
    if not rows:
        return '<div class="page"><p class="note">Shot Map V1.2 未找到</p></div>'

    if compact:
        cols = ["Shot", "SC", "功能", "P", "景别", "状态"]
    else:
        cols = ["Shot", "SC", "功能", "P", "景别", "机位", "视觉中心", "第二信息", "PNG ID", "页", "原文", "线索公平", "状态"]

    hdr = "".join(f"<th>{c}</th>" for c in cols)
    trs = ""
    for row in rows:
        tds = "".join(f"<td>{row.get(c, '')}</td>" for c in cols)
        trs += f"<tr>{tds}</tr>"

    cls = "gate-table shot-map" + (" compact" if compact else " full")
    return f"""
    <div class="page shot-map-page">
      <h2>分镜 Shot Map · A001（V1.2 完整分镜）</h2>
      <p class="lead">出典：<code>样章包/03_案01_分镜头与插页地图_V1.2_完整分镜.md</code></p>
      <table class="{cls}">
        <thead><tr>{hdr}</tr></thead>
        <tbody>{trs}</tbody>
      </table>
    </div>
    """


def shot_map_gallery_html() -> str:
    """Render depth_anchor images keyed by shot map, with placeholders for missing."""
    rows = parse_shot_map_v12()
    parts = ['<div class="page"><h2>Shot Map · 挿画対照（V1.2）</h2><p class="lead">L0 v0.2 · supplement v0.1 · 探索稿</p>']
    for row in rows:
        shot = row.get("Shot", "")
        png_id = row.get("PNG ID", "")
        func = row.get("功能", "")
        status = row.get("状态", "")
        if not shot or shot == "Shot":
            continue
        if png_id in ("线索卡_01_设计", "—", ""):
            parts.append(placeholder_box(shot, f"{func} · {status}"))
            continue
        candidates = resolve_depth_png(png_id)
        cap = f"{shot} · {func} · {status}"
        parts.append(figure_from_paths(candidates, "quarter", cap, shot))
    parts.append("</div>")
    return "\n".join(parts)


def inject_a001_shots(chapter: str, section: str | None) -> str:
    sec_num = None
    if section and section[0].isdigit():
        sec_num = section.split()[0]
    parts: list[str] = []
    for key, rule_sec, figs in SHOT_INJECT:
        if key not in chapter:
            continue
        if rule_sec is None and section is not None:
            continue
        if rule_sec is not None and sec_num != rule_sec:
            continue
        for paths, css, cap, sid in figs:
            parts.append(figure_from_paths(paths, css, cap, sid))
    return "\n".join(parts)


def prose_to_html_a001(text: str, inject_shots: bool = True) -> str:
    blocks: list[str] = []
    para: list[str] = []
    current_chapter = ""

    for line in text.splitlines():
        s = line.rstrip()
        if not s:
            if para:
                blocks.append("<p>" + "<br/>".join(para) + "</p>")
                para = []
            continue
        if s == "---":
            blocks.append('<hr class="scene-break"/>')
            continue
        if s.startswith("=") and len(s.strip()) > 10:
            continue
        if re.match(r"^(序|一、|めくれた|翘边|二、|三、|四、|五、)", s):
            current_chapter = s.split("-")[0].strip()
            blocks.append(f'<h2 class="chapter">{current_chapter}</h2>')
            if inject_shots:
                fig = inject_a001_shots(current_chapter, None)
                if fig:
                    blocks.append(fig)
            continue
        if s.startswith("### "):
            sec = s[4:]
            blocks.append(f'<h3 class="section">{sec}</h3>')
            if inject_shots:
                fig = inject_a001_shots(current_chapter, sec)
                if fig:
                    blocks.append(fig)
            continue
        if s.startswith("  ") and not s.startswith("  http"):
            para.append(f'<span class="indent">{s.strip()}</span>')
        elif re.match(r"^-{4,}$", s):
            continue
        else:
            clean = re.sub(r"\*\*(.+?)\*\*", r"<strong>\1</strong>", s)
            if clean.strip() and not clean.startswith("【"):
                para.append(clean)
    if para:
        blocks.append("<p>" + "<br/>".join(para) + "</p>")
    return "\n".join(blocks)


def four_column_review_html(scope: str, per_case: list[str] | None = None) -> str:
    """四栏审核：文字 / 科学 / 插画 / 日本文化."""
    cols = [
        ("文字", "E04 注音·文体·年级", "□ PASS □ HOLD □ 修正"),
        ("科学", "公平线索·DB1·实验安全", "□ PASS □ HOLD □ 修正"),
        ("插画", "V13 画风·角色一致·占位", "□ PASS □ HOLD □ 修正"),
        ("日本文化", "E07 侧廊·上履き·壁報", "□ PASS □ HOLD □ 修正"),
    ]
    grid = '<div class="review-grid">'
    for title, hint, checks in cols:
        grid += f"""
        <div class="review-col">
          <h3>{title}</h3>
          <p class="hint">{hint}</p>
          <p class="checks">{checks}</p>
          <p class="line">意见：_________________________</p>
          <p class="line">署名：__________ 日付：______</p>
        </div>"""
    grid += "</div>"

    case_blocks = ""
    if per_case:
        for cid in per_case:
            case_blocks += f'<h3 class="case-review">{cid}</h3>{grid}'

    return f"""
    <div class="page review-page">
      <h2>四栏审核 · {scope}</h2>
      <p class="lead">文字 / 科学 / 插画 / 日本文化 — 手写或打印填写</p>
      {case_blocks if case_blocks else grid}
    </div>
    """


def base_css() -> str:
    return """
@page { size: A5; margin: 14mm 11mm 18mm 14mm; }
* { box-sizing: border-box; }
body {
  font-family: "Yu Gothic", "Meiryo", "Hiragino Sans", sans-serif;
  font-size: 10.5pt; line-height: 1.58; color: #1a1a1a; max-width: 148mm; margin: 0 auto;
}
.page { page-break-after: always; padding: 3mm 0; min-height: 210mm; }
.page:last-child { page-break-after: auto; }
.cover {
  display:flex; flex-direction:column; justify-content:center; align-items:center;
  text-align:center; min-height:240mm;
  background: linear-gradient(165deg, #faf7f2 0%, #e8e0d4 100%);
}
.cover h1 { font-size: 22pt; margin: 0.3em 0; }
.cover .vol { font-size: 12pt; color: #444; }
.cover .badge {
  margin-top: 1em; font-size: 9pt; border: 2px solid #8a6d4b;
  padding: 0.4em 0.9em; color: #5a4030; background: #fffef8;
}
.cover img { max-width: 68%; margin: 0.8em 0; border: 1px solid #ccc; }
h2 { font-size: 12pt; border-bottom: 1px solid #bbb; padding-bottom: 0.2em; }
h3.section { font-size: 10pt; color: #555; }
.lead, .note { font-size: 8.5pt; color: #666; }
.gate-table { width:100%; border-collapse: collapse; font-size: 8pt; margin: 0.6em 0; }
.gate-table th, .gate-table td { border: 1px solid #ccc; padding: 0.28em 0.35em; vertical-align: top; }
.gate-table .st { font-weight: bold; white-space: nowrap; }
.shot-map.full { font-size: 6.5pt; }
.shot-map.compact { font-size: 8pt; }
.case-table { width:100%; border-collapse: collapse; font-size: 9pt; }
.case-table th, .case-table td { border: 1px solid #ccc; padding: 0.35em 0.45em; }
.case-strip {
  background: #f5f0e8; padding: 0.5em 0.8em; margin: 0.8em 0 0.4em;
  border-left: 4px solid #8a6d4b; font-weight: bold;
}
.illus { text-align:center; margin: 0.5em 0; page-break-inside: avoid; }
.illus img { max-width: 100%; height: auto; }
.illus.half img { max-height: 78mm; }
.illus.quarter img { max-height: 46mm; }
.illus figcaption { font-size: 7.5pt; color: #555; }
.shot-badge { color: #8a6d4b; font-size: 7pt; }
.placeholder-illus {
  border: 2px dashed #bbb; background: #f0f0f0; min-height: 36mm;
  display:flex; flex-direction:column; align-items:center; justify-content:center;
  color:#888; font-size:8.5pt; margin:0.5em 0; padding: 0.5em;
}
.placeholder-illus .ph-id { font-weight: bold; font-size: 10pt; color: #666; }
.placeholder-illus .ph-brief { font-size: 8pt; color: #999; margin-top: 0.2em; }
.diary-page {
  background: #fffef8; padding: 1em; border-left: 3px solid #d4c4a8;
  background-image: repeating-linear-gradient(#0000 0 1.3em, #e8e4dc 1.3em 1.31em);
}
.review-grid { display: grid; grid-template-columns: 1fr 1fr; gap: 0.6em; }
.review-col { border: 1px solid #ddd; padding: 0.5em; font-size: 8.5pt; }
.review-col h3 { font-size: 9.5pt; margin: 0 0 0.3em; border: none; }
.review-col .hint { color: #777; font-size: 7.5pt; }
.review-col .line { border-bottom: 1px dotted #ccc; min-height: 1.3em; margin: 0.4em 0; }
.clue-box {
  font-family: "Yu Gothic", monospace; font-size: 9pt;
  background: #f7f4ef; padding: 1em; border: 1px solid #ccc; white-space: pre-wrap;
}
.appendix-cn { background: #fafafa; }
.appendix-cn .prose-cn { font-family: "Microsoft YaHei", sans-serif; }
blockquote.tagline { font-size: 14pt; text-align:center; border:none; color:#6a5040; }
p { margin: 0.45em 0; }
p .indent { display:block; padding-left:1.2em; }
hr.scene-break { border:none; border-top:1px dotted #ccc; margin:0.8em 0; }
.toc { list-style:none; padding:0; }
.toc li { padding: 0.3em 0; border-bottom: 1px dotted #ccc; }
@media print { body { max-width: none; } }
"""


def gate_rows_v01a() -> list[tuple[str, str, str, str]]:
    return [
        ("正典", "第一話=A001 非湿椅子", "✅ LOCK", "序+めくれたポスター"),
        ("正典", "Vol1=Plan B", "✅ LOCK", "Gate A · 本包は事件①のみ"),
        ("Gate A", "P0-01/02 IP LOCK", "✅", "2026-06-05"),
        ("C03", "公平线索≥3", "✅", "風側/雨天/貼り方"),
        ("C04", "温柔真相", "✅", "志郎非悪意"),
        ("V13", "桥梁书年龄感", "🟡", "v0.1/v0.2 探索稿"),
        ("E04", "日文文体", "🟡", "JP_VOICE_v1.1 · 田中待ち"),
        ("E07", "校园五维", "⬜", "側廊/クラブ/壁報"),
        ("科学", "P0-04 A001", "⬜", "DB1+DA3 待签"),
        ("S20", "非湿椅子叙事", "✅", "实验=テープ+ドライヤー"),
    ]


def gate_rows_v01b() -> list[tuple[str, str, str, str]]:
    return [
        ("正典", "Vol1=Plan B A001–A005", "✅ LOCK", "非湿椅子 C001"),
        ("正典", "首市场=日本", "✅ LOCK", "学堂奇事録"),
        ("正典", "A001=めくれたポスター", "✅ LOCK", "深度锚点 v0.2"),
        ("C03", "5案轻推理结构", "✅", "L1 尾钩串联"),
        ("C03b", "公平线索", "✅/🟡", "A001–A003 定稿 · A004–A005 初稿"),
        ("V13", "桥梁书年龄感", "🟡", "探索稿·非印刷"),
        ("E04", "日文文体", "🟡", "A001 JP_VOICE_v1.1 · ②–⑤初译"),
        ("E07", "田中校园五维", "⬜", "待全文督查"),
        ("E20", "出版社提案", "🟡", "本包可发对外"),
        ("S20", "红线", "✅", "无湿椅子主叙事"),
    ]


def gate_table_html(rows: list[tuple[str, str, str, str]], title: str, lead: str) -> str:
    trs = "".join(
        f"<tr><td>{a}</td><td>{b}</td><td class='st'>{c}</td><td>{d}</td></tr>" for a, b, c, d in rows
    )
    return f"""
    <div class="page">
      <h2>{title}</h2>
      <p class="lead">{lead}</p>
      <table class="gate-table">
        <thead><tr><th>項目</th><th>区分</th><th>状態</th><th>備考</th></tr></thead>
        <tbody>{trs}</tbody>
      </table>
      <p class="note">⬜=待人审 · 🟡=可发探索稿 · ✅=正典锁定</p>
    </div>
    """


def positioning_html() -> str:
    return """
    <div class="page">
      <h2>IP定位 · 出版提案摘要</h2>
      <blockquote class="tagline">謎が、解ける日。</blockquote>
      <p><strong>対象</strong>：7–11歳（主 8–10）· 校园軽推理成長 IP</p>
      <p><strong>設計原則</strong>：かわいさで入る · 謎で留まる · キャラで戻る · 観察で続く</p>
      <p><strong>境界</strong>：硬科普/教辅/恐怖/大人破案 ✗ · 日常小異常+温柔真相 ✓</p>
      <p><strong>Vol1</strong>：Plan B · 5年2組 陸珣（しゅん）視点 · 観察クラブ入門 · 5事件合订</p>
      <p class="src">Gate A：NOT 湿椅子 · A001=めくれたポスター · IP LOCK 2026-06-05</p>
    </div>
    """


def case_overview_html() -> str:
    rows = "".join(
        f"<tr><td>{a}</td><td>{n}</td><td>{t}</td><td>{k}</td><td>{s}</td></tr>"
        for a, n, t, k, _fname, s in CASES
    )
    return f"""
    <div class="page">
      <h2>事件一覧 · Case Card 要約（A001–A005）</h2>
      <table class="case-table">
        <thead><tr><th>ID</th><th>案</th><th>タイトル</th><th>核心</th><th>稿</th></tr></thead>
        <tbody>{rows}</tbody>
      </table>
      <p class="note">L1 尾钩：空欄→投稿→裏の路→中谷半句→瑆誤描 · L3 巻末壁報第2起</p>
    </div>
    """


def hikaru_page_html() -> str:
    img = first_existing([
        ILLUST_SAMPLE / "V-S02_瑆日记页_v0.1.png",
        ILLUST / "V-S02_陸瑆日记页.png",
    ])
    img_html = figure_from_paths([img], "half", "V-S02 · 瑆のノート（B轨）", "V-S02") if img else ""
    body = bde.strip_txt_meta(bde.read_text(HIKARU_TXT)) if HIKARU_TXT.exists() else ""
    text_html = prose_to_html_a001(body, inject_shots=False) if body else ""
    return f"""
    <div class="page diary-page">
      <h2>陸瑆（ひかる）のノート · 事件①</h2>
      {img_html}
      {text_html}
      <p class="note">B轨：正文层と画风を区別 · 問句收束 · 不剧透</p>
    </div>
    """


def db1_and_experiment_html() -> str:
    db1 = figure_from_paths(
        [DEPTH / "V-S01-B1_风侧机制图_v0.2.png", DEPTH / "V-S01-B1_风侧机制图.png"],
        "half", "DB1 · 風側メカニズム SUM", "DB1",
    )
    exp = bde.read_text(EXPERIMENT) if EXPERIMENT.exists() else ""
    exp_html = prose_to_html_a001(re.sub(r"^# .+\n", "", exp), inject_shots=False)
    return f"""
    <div class="page">
      <h2>家庭実験 · 事件①</h2>
      <p class="lead">カード紙+テープ+ドライヤー · 大人同伴 · 非湿椅子实验</p>
      {db1}
      <div class="experiment-text">{exp_html}</div>
    </div>
    """


def clue_card_html() -> str:
    raw = bde.read_text(CLUE_MD) if CLUE_MD.exists() else ""
    pre = re.search(r"```\n([\s\S]*?)```", raw)
    box = pre.group(1) if pre else "线索卡 #01 · CL-01 · 见 07_线索卡_设计稿.md"
    return f"""
    <div class="page clue-card">
      <h2>観察クラブ · 手がかりカード #01（CL-01）</h2>
      <pre class="clue-box">{box.strip()}</pre>
      <p class="note">三项线索正文均已出现 · 仅回顾 · ⬜ 版式 E22 待决</p>
    </div>
    """


def cn_appendix_html() -> str:
    raw = bde.strip_txt_meta(bde.read_text(CN_A001)) if CN_A001.exists() else ""
    body = prose_to_html_a001(raw, inject_shots=False)
    return f"""
    <div class="page appendix-cn">
      <h2>中文附录 · 对照用（非主交付语言）</h2>
      <p class="lead">源：<code>01_正文/案01_翘边的海报_HybridVoice_V1.1定稿.txt</code></p>
      <p class="note">正典备注：Vol1=Plan B · A001 · 非湿椅子 · 珣=しゅん / 瑆=ひかる</p>
      <div class="prose-cn">{body}</div>
    </div>
    """


def build_v01a_html() -> str:
    jp_src = JP_A001 if JP_A001.exists() else JP_SAMPLE
    jp_raw = bde.strip_txt_meta(bde.read_text(jp_src))
    body = prose_to_html_a001(jp_raw)
    cover_img = first_existing([DEPTH / "V-S01-A1_侧廊发现_v0.2.png", ILLUST / "V-S01_侧廊海报.png"])
    cover_uri = bde.img_data_uri(cover_img) if cover_img else ""

    chars = (
        "【事件① 出场】\n"
        "陸珣（しゅん）· 5年2組 · 転校生\n"
        "伊藤光 · 5年2組 · おもしろ観察クラブ\n"
        "加藤慧美 · 5年1組 · 記録\n"
        "松本志郎 · 5年3組 · 検証\n"
        "陸瑆（ひかる）· 4年2組 · 日記层（B轨）"
    )

    return f"""<!DOCTYPE html>
<html lang="ja">
<head>
<meta charset="utf-8"/>
<title>学堂奇事録 V0.1-A 第一話审阅包</title>
<style>{base_css()}</style>
</head>
<body>

<div class="page cover">
  <p class="badge">V0.1-A · 第一話审阅包 · 非売品 · RGB探索稿</p>
  <h1>学堂奇事録</h1>
  <p class="vol">第1巻 · 事件①「めくれたポスター」<br/>序 + A001 · Plan B · Gate A</p>
  {"<img src='" + cover_uri + "' alt='DA1'/>" if cover_uri else ""}
  <p class="note">{OUT_DATE} · 日本語主 · 中文附录</p>
</div>

<div class="page">
  <h2>出版声明</h2>
  <ul style="font-size:9.5pt;line-height:1.55;">
    <li><strong>種別</strong>：V0.1-A 团队共审里程碑（专家/编辑/美工）</li>
    <li><strong>非</strong>：印厂原稿 · CMYK终稿 · 湿椅子(C001)卷</li>
    <li><strong>正典</strong>：Vol1 Plan B · A001=めくれたポスター · IP LOCK</li>
  </ul>
</div>

{gate_table_html(gate_rows_v01a(), "正典 · V0.1-A 门禁对照", "第一話里程碑 · 序 + A001")}
{positioning_html()}
{shot_map_table_html(compact=False)}

<div class="page">
  <h2>登場人物（事件①）</h2>
  <pre style="white-space:pre-wrap;font-family:inherit;font-size:10pt;">{chars}</pre>
</div>

<div class="page">
  <h2>目次 · V0.1-A</h2>
  <ul class="toc">
    <li>序 · 四月の第二月曜日</li>
    <li>一、めくれたポスター（§1–4 + depth_anchor）</li>
    <li>陸瑆のノート · V-S02</li>
    <li>手がかりカード CL-01</li>
    <li>DB1 + 家庭実験</li>
    <li>Shot Map V1.2 挿画対照</li>
    <li>四栏审核（A001）</li>
    <li>中文附录</li>
  </ul>
</div>

<div class="page prose">{body}</div>

{hikaru_page_html()}
{clue_card_html()}
{db1_and_experiment_html()}
{shot_map_gallery_html()}
{four_column_review_html("V0.1-A · A001")}
{cn_appendix_html()}

<div class="page">
  <h2>版本记录 · V0.1-A</h2>
  <table class="gate-table">
    <tr><th>版</th><td>V0.1-A · 第一話审阅包</td></tr>
    <tr><th>日期</th><td>{OUT_DATE}</td></tr>
    <tr><th>范围</th><td>序 + A001 · 非全卷 · 非湿椅子</td></tr>
    <tr><th>主语言</th><td>日本語 + 中文附录</td></tr>
    <tr><th>插图</th><td>L0 v0.2 + supplement v0.1 · 探索稿·非印刷</td></tr>
  </table>
  <p class="note">rebuild: python 05_出版成果/tools/build_expert_publication_pack.py --mode v01a</p>
</div>

</body>
</html>
"""


def case_art_html(case_id: str) -> str:
    if case_id == "A001":
        parts = []
        for _key, _sec, figs in SHOT_INJECT:
            for paths, css, cap, sid in figs:
                parts.append(figure_from_paths(paths, css, cap, sid))
        db1 = figure_from_paths(
            [DEPTH / "V-S01-B1_风侧机制图_v0.2.png", DEPTH / "V-S01-B1_风侧机制图.png"],
            "quarter", "DB1 · 机制SUM", "DB1",
        )
        parts.append(db1)
        return "\n".join(parts)

    parts = []
    for fname, cap in CASE_ART.get(case_id, []):
        path = ILLUST / fname
        if path.exists():
            parts.append(figure_from_paths([path], "quarter", cap, fname.split("_")[0]))
        else:
            parts.append(placeholder_box(fname.split("_")[0], cap))
    if not parts:
        parts.append(placeholder_box(case_id, "挿画プレースホルダ · 探索稿"))
    return "\n".join(parts)


def load_vol1_jp_text() -> str:
    bde.ensure_jp_canon()
    if CANON_JP.exists():
        return bde.strip_txt_meta(bde.read_text(CANON_JP))
    chunks: list[str] = []
    for _cid, _n, _t, _k, fname, _s in CASES:
        p = BODY_DIR / fname
        if p.exists():
            chunks.append(bde.strip_txt_meta(bde.read_text(p)))
    return "\n\n".join(chunks)


def build_v01b_html() -> str:
    canon_raw = load_vol1_jp_text()
    body_html = bde.md_prose_to_html(canon_raw, "ja")
    notes_html = bde.notes_to_html(bde.read_text(NOTES_JP), "ja") if NOTES_JP.exists() else ""

    # Per-case sections with title strips + art
    case_sections = ""
    for cid, num, title, core, fname, status in CASES:
        p = BODY_DIR / fname
        text = bde.strip_txt_meta(bde.read_text(p)) if p.exists() else f"（{fname} 未找到）"
        if cid == "A001":
            case_body = prose_to_html_a001(text)
        else:
            case_body = bde.md_prose_to_html(text, "ja")
        case_sections += f"""
        <div class="page case-page">
          <div class="case-strip">{cid} {num} · {title} · {core} · {status}</div>
          <p class="note">探索稿·非印刷 — 挿画は v0.9 プレースホルダ</p>
          {case_art_html(cid)}
          <div class="prose case-text">{case_body}</div>
        </div>
        """

    cover_img = first_existing([DEPTH / "V-S01-A1_侧廊发现_v0.2.png", ILLUST / "V-S01_侧廊海报.png"])
    cover_uri = bde.img_data_uri(cover_img) if cover_img else ""

    return f"""<!DOCTYPE html>
<html lang="ja">
<head>
<meta charset="utf-8"/>
<title>学堂奇事録 V0.1-B 全卷提案包</title>
<style>{base_css()}</style>
</head>
<body>

<div class="page cover">
  <p class="badge">V0.1-B · Vol1 全卷提案包 · 出版社对外 · RGB</p>
  <h1>学堂奇事録</h1>
  <p class="vol">第1巻 · おかしいと思ったら、まず見てみる<br/>5事件骨架 · Plan B</p>
  {"<img src='" + cover_uri + "' alt='cover'/>" if cover_uri else ""}
  <p class="note">{OUT_DATE} · 探索稿·非印刷 · Gate A 正典</p>
</div>

<div class="page">
  <h2>出版声明 · V0.1-B</h2>
  <ul style="font-size:9.5pt;line-height:1.55;">
    <li><strong>種別</strong>：Vol1 全卷提案包（出版社/专家对外骨架）</li>
    <li><strong>正典</strong>：Plan B · A001=めくれたポスター · NOT 湿椅子</li>
    <li><strong>插图</strong>：A001=depth_anchor v0.2 · 他案=探索稿占位</li>
  </ul>
</div>

{gate_table_html(gate_rows_v01b(), "正典 · V0.1-B 门禁对照", "全卷提案 · 5案骨架")}
{positioning_html()}
{case_overview_html()}
{shot_map_table_html(compact=True)}

<div class="page">
  <h2>目次 · V0.1-B</h2>
  <ul class="toc">
    <li>序 + 事件①–⑤ 全文（日本語）</li>
    <li>陸瑆（ひかる）のノート · Vol1</li>
    <li>事件別挿画対照（探索稿）</li>
    <li>四栏审核（全卷）</li>
  </ul>
</div>

<div class="page prose">
  <h2>全文日本語 · 序 + 事件①–⑤</h2>
  <p class="note"> assembled正本 · 挿画は下記事件別セクション参照</p>
  {body_html}
</div>

{case_sections}

<div class="page diary-page">
  <h2>陸瑆（ひかる）のノート · Vol1</h2>
  {notes_html}
</div>

{four_column_review_html("V0.1-B · 全卷", per_case=[c[0] for c in CASES])}

<div class="page">
  <h2>版本记录 · V0.1-B</h2>
  <table class="gate-table">
    <tr><th>版</th><td>V0.1-B · Vol1 全卷提案包</td></tr>
    <tr><th>日期</th><td>{OUT_DATE}</td></tr>
    <tr><th>范围</th><td>序 + A001–A005 · 5案骨架</td></tr>
    <tr><th>插图</th><td>探索稿·非印刷 · A001 depth_anchor 実装</td></tr>
  </table>
  <p class="note">rebuild: python 05_出版成果/tools/build_expert_publication_pack.py --mode v01b</p>
</div>

</body>
</html>
"""


def write_gate_md(rows: list[tuple[str, str, str, str]], title: str, script_note: str, filename: str) -> Path:
    lines = [f"# {title} · {OUT_DATE}", "", "| 项 | 类别 | 状态 | 备注 |", "|----|------|------|------|"]
    for a, b, c, d in rows:
        lines.append(f"| {a} | {b} | {c} | {d} |")
    lines.extend(["", f"生成：`{script_note}`", ""])
    path = OUT / filename
    path.write_text("\n".join(lines), encoding="utf-8")
    return path


def _ppt_setup():
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        subprocess.run([sys.executable, "-m", "pip", "install", "python-pptx", "-q"], check=True)
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
    return Presentation, Inches, Pt, PP_ALIGN


def build_v01a_ppt(pdf_hint: Path) -> Path:
    Presentation, Inches, Pt, PP_ALIGN = _ppt_setup()
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    blank = prs.slide_layouts[6]
    MX = Inches(0.65)

    def title(t: str, sub: str = "") -> None:
        s = prs.slides.add_slide(blank)
        box = s.shapes.add_textbox(MX, Inches(2.2), Inches(14.7), Inches(4))
        tf = box.text_frame
        tf.text = t
        tf.paragraphs[0].font.size = Pt(36)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        if sub:
            p = tf.add_paragraph()
            p.text = sub
            p.font.size = Pt(18)
            p.alignment = PP_ALIGN.CENTER

    def bullets(t: str, lines: list[str]) -> None:
        s = prs.slides.add_slide(blank)
        tb = s.shapes.add_textbox(MX, Inches(0.5), Inches(14.7), Inches(8))
        tf = tb.text_frame
        tf.text = t
        tf.paragraphs[0].font.size = Pt(28)
        tf.paragraphs[0].font.bold = True
        for line in lines:
            p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(16)
            p.space_before = Pt(8)

    title("V0.1-A · 第一話审阅包", "序 + A001 めくれたポスター · 日本語主 + 中文附录")
    bullets("Gate A 正典", [
        "• NOT 湿椅子 · Plan B · A001=めくれたポスター",
        "• 专家/编辑/美工并行审稿",
        "• 插图=探索稿 v0.1/v0.2 · RGB",
        f"• PDF：{pdf_hint.name}",
    ])
    bullets("Shot Map V1.2", [
        "L0: DA1–DA4 + DB1 + DC1（v0.2）",
        "L1: V-S02 瑆 + DA2b/DA3b/DA4b（v0.1）",
        "L2: V-S00 序钩 · CL-01 线索卡",
        "四栏审核：文字/科学/插画/日本文化",
    ])
    img = first_existing([DEPTH / "V-S01-A1_侧廊发现_v0.2.png"])
    if img:
        s = prs.slides.add_slide(blank)
        s.shapes.add_textbox(MX, Inches(0.4), Inches(6), Inches(0.5)).text_frame.text = "DA1 · depth_anchor v0.2"
        s.shapes.add_picture(str(img), Inches(5.5), Inches(0.6), width=Inches(10))

    out = OUT / V01A["ppt"]
    prs.save(str(out))
    return out


def build_v01b_ppt(pdf_hint: Path) -> Path:
    Presentation, Inches, Pt, PP_ALIGN = _ppt_setup()
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    blank = prs.slide_layouts[6]
    MX = Inches(0.65)
    CW = Inches(14.7)

    def title(t: str, sub: str = "") -> None:
        s = prs.slides.add_slide(blank)
        box = s.shapes.add_textbox(MX, Inches(2.0), CW, Inches(4.5))
        tf = box.text_frame
        tf.text = t
        tf.paragraphs[0].font.size = Pt(40)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        if sub:
            p = tf.add_paragraph()
            p.text = sub
            p.font.size = Pt(18)
            p.alignment = PP_ALIGN.CENTER

    def bullets(t: str, lines: list[str], cols: int = 1) -> None:
        s = prs.slides.add_slide(blank)
        title_box = s.shapes.add_textbox(MX, Inches(0.45), CW, Inches(0.7))
        title_box.text_frame.text = t
        title_box.text_frame.paragraphs[0].font.size = Pt(30)
        title_box.text_frame.paragraphs[0].font.bold = True
        col_w = (CW - Inches(0.4)) / cols if cols > 1 else CW
        for ci in range(cols):
            chunk = lines[ci :: cols] if cols > 1 else lines
            left = MX + ci * (col_w + Inches(0.4))
            tb = s.shapes.add_textbox(left, Inches(1.15), col_w, Inches(7.2))
            tf = tb.text_frame
            tf.word_wrap = True
            for i, line in enumerate(chunk):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = line
                p.font.size = Pt(17 if cols == 1 else 15)
                p.space_after = Pt(6)

    title("V0.1-B · Vol1 全卷提案包", "5案骨架 · 出版社对外 · 探索稿·非印刷")
    bullets("IP定位", [
        "• 7–11歳 · 名古屋 · 校园軽推理",
        "• Gate A：NOT 湿椅子 · Plan B",
        "• かわいさで入る · 謎で留まる · 観察で続く",
        f"• PDF：{pdf_hint.name}",
    ])
    case_lines = [f"{a} {n} {t} · {k} {s}" for a, n, t, k, _f, s in CASES]
    bullets("事件 A001–A005", case_lines, cols=2)
    bullets("成果物構成", [
        "序 + 5事件全文（日本語）",
        "陸瑆笔记 Vol1",
        "A001 depth_anchor · 他案占位挿画",
        "四栏审核 × 5案",
        "门禁对照表 MD",
    ])
    img = first_existing([DEPTH / "V-S01-A1_侧廊发现_v0.2.png", ILLUST / "V-S01_侧廊海报.png"])
    if img:
        s = prs.slides.add_slide(blank)
        s.shapes.add_textbox(MX, Inches(0.4), Inches(5.5), Inches(0.6)).text_frame.text = "主视觉 · A001 depth_anchor"
        s.shapes.add_picture(str(img), Inches(5.8), Inches(0.7), width=Inches(9.5))

    out = OUT / V01B["ppt"]
    prs.save(str(out))
    return out


def run_prep() -> None:
    if ILLUST_TOOLS.joinpath("generate_vol1_illustrations.py").exists():
        subprocess.run([sys.executable, str(ILLUST_TOOLS / "generate_vol1_illustrations.py")], check=True)
    subprocess.run([sys.executable, str(DISPLAY_TOOLS / "assemble_jp_canon.py")], check=True)


def build_pack(mode: str) -> int:
    OUT.mkdir(parents=True, exist_ok=True)
    rc = 0

    if mode in ("v01a", "both"):
        html = build_v01a_html()
        html_path = OUT / V01A["html"]
        html_path.write_text(html, encoding="utf-8")
        print(f"[V0.1-A] HTML: {html_path}")

        pdf_path = OUT / V01A["pdf"]
        if not bde.html_to_pdf_chrome(html_path, pdf_path):
            print("[V0.1-A] PDF failed — open HTML and print", file=sys.stderr)
            rc = 1
        else:
            print(f"[V0.1-A] PDF: {pdf_path} ({pdf_path.stat().st_size // 1024} KB)")

        gate = write_gate_md(
            gate_rows_v01a(),
            "V0.1-A 第一話审阅包 · 门禁对照",
            "build_expert_publication_pack.py --mode v01a",
            V01A["gate"],
        )
        print(f"[V0.1-A] Gate: {gate}")

        if rc == 0:
            ppt = build_v01a_ppt(pdf_path)
            print(f"[V0.1-A] PPT: {ppt} ({ppt.stat().st_size // 1024} KB)")

    if mode in ("v01b", "both"):
        html = build_v01b_html()
        html_path = OUT / V01B["html"]
        html_path.write_text(html, encoding="utf-8")
        print(f"[V0.1-B] HTML: {html_path}")

        pdf_path = OUT / V01B["pdf"]
        if not bde.html_to_pdf_chrome(html_path, pdf_path):
            print("[V0.1-B] PDF failed — open HTML and print", file=sys.stderr)
            rc = 1
        else:
            print(f"[V0.1-B] PDF: {pdf_path} ({pdf_path.stat().st_size // 1024} KB)")

        gate = write_gate_md(
            gate_rows_v01b(),
            "V0.1-B 全卷提案包 · 门禁对照",
            "build_expert_publication_pack.py --mode v01b",
            V01B["gate"],
        )
        print(f"[V0.1-B] Gate: {gate}")

        if rc == 0:
            ppt = build_v01b_ppt(pdf_path)
            print(f"[V0.1-B] PPT: {ppt} ({ppt.stat().st_size // 1024} KB)")

    return rc


def main() -> int:
    parser = argparse.ArgumentParser(description="Build Vol1 V0.1-A / V0.1-B publication packs")
    parser.add_argument(
        "--mode",
        choices=("v01a", "v01b", "both"),
        default="both",
        help="v01a=第一話审阅包 · v01b=全卷提案包 · both=default",
    )
    args = parser.parse_args()

    if args.mode in ("v01a", "both") and not JP_A001.exists() and not JP_SAMPLE.exists():
        print("Missing A001 JP source", file=sys.stderr)
        return 1

    run_prep()
    return build_pack(args.mode)


if __name__ == "__main__":
    sys.exit(main())
